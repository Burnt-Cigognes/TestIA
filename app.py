import streamlit as st
from openai import OpenAI
import PyPDF2
from pptx import Presentation
import base64
import io
from datetime import datetime
import json
import requests
from urllib.parse import quote

# ==========================================
# 1. CONFIGURATION
# ==========================================
st.set_page_config(
    page_title="IA CIC",
    page_icon="🙉",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ==========================================
# 2. CSS
# ==========================================
st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Cormorant+Garamond:wght@400;500;600&family=DM+Sans:opsz,wght@9..40,300;9..40,400;9..40,500&display=swap');

    html, body, [class*="css"] { font-family: 'DM Sans', sans-serif; }
    .stApp { background: #10131c; color: #ddd8ce; }
    #MainMenu, footer, header { visibility: hidden; }

    /* Sidebar */
    section[data-testid="stSidebar"] > div { background: #131722; padding-top: 0; }
    [data-testid="stSidebar"] { border-right: 1px solid #1f2535; }
    [data-testid="stSidebar"] label {
        color: #6b7089 !important;
        font-size: 0.7rem !important;
        letter-spacing: 0.12em !important;
        text-transform: uppercase !important;
        font-weight: 500 !important;
    }

    /* Selectbox */
    [data-testid="stSelectbox"] > div > div {
        background: #1b2030 !important;
        border: 1px solid #252c3f !important;
        border-radius: 8px !important;
        color: #ddd8ce !important;
        font-size: 0.85rem !important;
    }
    [data-testid="stSelectbox"] svg { fill: #6b7089; }

    /* Text area */
    .stTextArea textarea {
        background: #1b2030 !important;
        border: 1px solid #252c3f !important;
        border-radius: 8px !important;
        color: #ddd8ce !important;
        font-size: 0.82rem !important;
        font-family: 'DM Sans', sans-serif !important;
    }
    .stTextArea textarea:focus {
        border-color: #b8995a !important;
        box-shadow: 0 0 0 2px rgba(184,153,90,0.15) !important;
    }

    /* Toggle */
    [data-testid="stToggle"] label {
        color: #ddd8ce !important;
        font-size: 0.82rem !important;
        text-transform: none !important;
        letter-spacing: 0 !important;
    }

    /* File uploader */
    [data-testid="stFileUploader"] section {
        background: #1b2030 !important;
        border: 1px solid #252c3f !important;
        border-radius: 8px !important;
        padding: 12px !important;
    }
    [data-testid="stFileUploader"] section:hover { border-color: #b8995a !important; }
    [data-testid="stFileUploader"] span,
    [data-testid="stFileUploader"] p,
    [data-testid="stFileUploader"] small,
    [data-testid="stFileUploader"] div {
        color: #ddd8ce !important;
        font-family: 'DM Sans', sans-serif !important;
        font-size: 0.82rem !important;
    }
    [data-testid="stFileUploader"] button {
        background: #1b2030 !important;
        border: 1px solid #252c3f !important;
        color: #ddd8ce !important;
        border-radius: 6px !important;
        font-size: 0.8rem !important;
        width: auto !important;
    }
    [data-testid="stFileUploader"] button:hover { border-color: #b8995a !important; color: #b8995a !important; }

    /* Buttons */
    .stButton > button {
        background: transparent !important;
        border: 1px solid #252c3f !important;
        color: #6b7089 !important;
        border-radius: 7px !important;
        font-family: 'DM Sans', sans-serif !important;
        font-size: 0.8rem !important;
        transition: all 0.2s ease !important;
        width: 100%;
    }
    .stButton > button:hover { border-color: #b8995a !important; color: #b8995a !important; background: rgba(184,153,90,0.06) !important; }

    /* Download button */
    [data-testid="stDownloadButton"] button {
        background: rgba(184,153,90,0.1) !important;
        border: 1px solid #b8995a !important;
        color: #b8995a !important;
        border-radius: 7px !important;
        font-size: 0.8rem !important;
        width: 100%;
        transition: all 0.2s ease !important;
    }
    [data-testid="stDownloadButton"] button:hover { background: rgba(184,153,90,0.2) !important; }

    /* Chat messages */
    [data-testid="stChatMessage"] {
        background: transparent;
        padding: 12px 0;
        border-bottom: 1px solid #191e2b;
        color: #ddd8ce !important;
    }
    [data-testid="stChatMessage"]:last-child { border-bottom: none; }
    [data-testid="stChatMessage"] p,
    [data-testid="stChatMessage"] li,
    [data-testid="stChatMessage"] span,
    [data-testid="stChatMessage"] div { color: #ddd8ce !important; }
    [data-testid="stChatMessage"] strong { color: #f0ebe0 !important; }
    [data-testid="stChatMessage"] h1,
    [data-testid="stChatMessage"] h2,
    [data-testid="stChatMessage"] h3 { color: #f0ebe0 !important; }
    [data-testid="stChatMessage"] code { color: #b8995a !important; background: #1b2030 !important; }

    /* Zone fixe bas */
    [data-testid="stBottom"],
    [data-testid="stBottom"] > div {
        background: #10131c !important;
        border-top: 1px solid #1f2535 !important;
    }

    /* Champ de saisie — blanc casse, texte fonce */
    [data-testid="stChatInput"] {
        background: #f5f2ec !important;
        border: 1px solid #d6cfc3 !important;
        border-radius: 12px !important;
        box-shadow: 0 2px 12px rgba(0,0,0,0.25) !important;
    }
    [data-testid="stChatInput"] textarea {
        background: #f5f2ec !important;
        color: #1a1e2a !important;
        font-family: 'DM Sans', sans-serif !important;
        font-size: 0.92rem !important;
    }
    [data-testid="stChatInput"] textarea::placeholder { color: #8a8070 !important; }
    [data-testid="stChatInput"]:focus-within {
        border-color: #b8995a !important;
        box-shadow: 0 0 0 3px rgba(184,153,90,0.18), 0 2px 12px rgba(0,0,0,0.25) !important;
    }
    [data-testid="stChatInput"] button { background: #b8995a !important; border: none !important; border-radius: 8px !important; color: #10131c !important; }
    [data-testid="stChatInput"] button:hover { background: #caa96a !important; }

    /* Tabs */
    .stTabs [data-baseweb="tab-list"] {
        background: transparent;
        gap: 4px;
        border-bottom: 1px solid #1f2535;
    }
    .stTabs [data-baseweb="tab"] {
        background: transparent !important;
        color: #6b7089 !important;
        border: none !important;
        border-bottom: 2px solid transparent !important;
        font-family: 'DM Sans', sans-serif !important;
        font-size: 0.82rem !important;
        padding: 8px 16px !important;
        letter-spacing: 0.04em;
    }
    .stTabs [aria-selected="true"] { color: #b8995a !important; border-bottom: 2px solid #b8995a !important; }

    /* Scrollbar */
    ::-webkit-scrollbar { width: 4px; }
    ::-webkit-scrollbar-track { background: #10131c; }
    ::-webkit-scrollbar-thumb { background: #252c3f; border-radius: 4px; }
    ::-webkit-scrollbar-thumb:hover { background: #b8995a; }

    /* HTML custom */
    .sidebar-header { padding: 24px 16px 16px; border-bottom: 1px solid #1f2535; margin-bottom: 16px; }
    .sidebar-brand { font-family: 'Cormorant Garamond', serif; font-size: 1.35rem; font-weight: 600; color: #ddd8ce; }
    .sidebar-brand span { color: #b8995a; }
    .sidebar-tagline { font-size: 0.68rem; color: #6b7089; letter-spacing: 0.14em; text-transform: uppercase; margin-top: 4px; }
    .section-divider { height: 1px; background: #1f2535; margin: 16px 0; }
    .section-label { font-size: 0.67rem; color: #6b7089; letter-spacing: 0.14em; text-transform: uppercase; font-weight: 500; margin-bottom: 10px; }
    .file-chip { display: inline-flex; align-items: center; gap: 6px; background: #1b2030; border: 1px solid #252c3f; border-radius: 20px; padding: 4px 10px; font-size: 0.75rem; color: #6b7089; margin: 6px 0; max-width: 100%; overflow: hidden; text-overflow: ellipsis; white-space: nowrap; }
    .chat-header { padding: 28px 0 16px; border-bottom: 1px solid #1f2535; margin-bottom: 8px; }
    .chat-title { font-family: 'Cormorant Garamond', serif; font-size: 1.6rem; font-weight: 500; color: #ddd8ce; }
    .chat-title span { color: #b8995a; }
    .chat-meta { font-size: 0.73rem; color: #6b7089; margin-top: 4px; }
    .welcome-box { background: linear-gradient(135deg, #161b28, #1a1f2e); border: 1px solid #1f2535; border-left: 3px solid #b8995a; border-radius: 10px; padding: 24px 28px; margin: 20px 0; }
    .welcome-box h3 { font-family: 'Cormorant Garamond', serif; font-size: 1.2rem; font-weight: 500; color: #ddd8ce; margin: 0 0 10px; }
    .welcome-box p { font-size: 0.85rem; color: #8b90a8; margin: 0; line-height: 1.7; }
    .capability-grid { display: grid; grid-template-columns: 1fr 1fr; gap: 10px; margin-top: 16px; }
    .capability-item { background: #10131c; border: 1px solid #1f2535; border-radius: 8px; padding: 12px 14px; font-size: 0.8rem; color: #8b90a8; }
    .capability-item .icon { font-size: 1rem; margin-bottom: 4px; display: block; }
    .stat-row { display: flex; gap: 8px; margin-top: 12px; flex-wrap: wrap; }
    .stat-pill { background: #10131c; border: 1px solid #1f2535; border-radius: 20px; padding: 3px 10px; font-size: 0.72rem; color: #6b7089; }
</style>
""", unsafe_allow_html=True)


# ==========================================
# 3. MODELES CHAT
# ==========================================
MODELS = {
    "💡 Gemini 2.0 Flash Lite — Economique": {
        "id": "google/gemini-2.0-flash-lite-001",
        "tier": "eco",
        "price": "~0,01 EUR / 10 000 mots",
        "desc": "Mails, actualites, resumes, questions simples",
    },
    "⚡ Gemini 2.5 Flash — Standard": {
        "id": "google/gemini-2.5-flash",
        "tier": "standard",
        "price": "~0,06 EUR / 10 000 mots",
        "desc": "Lecture PDF, etudes de marche, recherche acquereurs",
    },
    "🧠 Claude Sonnet 4.5 — Premium": {
        "id": "anthropic/claude-sonnet-4-5",
        "tier": "premium",
        "price": "~0,25 EUR / 10 000 mots",
        "desc": "Correction PPT/Word, analyses complexes, gros fichiers",
    },
}

TIER_STYLES = {
    "eco":      {"label": "ECONOMIQUE", "color": "#5a9e7a"},
    "standard": {"label": "STANDARD",  "color": "#b8995a"},
    "premium":  {"label": "PREMIUM",   "color": "#c06060"},
}

# Styles image generation
IMAGE_STYLES = {
    "Photoréaliste": "photorealistic, ultra detailed, 8k, professional photography",
    "Illustration": "digital illustration, clean lines, vibrant colors, professional",
    "Minimaliste": "minimalist, clean, simple, elegant, flat design",
    "Cinématique": "cinematic, dramatic lighting, film still, wide angle",
    "Aquarelle": "watercolor painting, soft colors, artistic, hand painted",
}

DEFAULT_SYSTEM = (
    "Tu es un assistant IA specialise en fusions-acquisitions (M&A), finance d entreprise et analyse strategique. "
    "Tu assistes un analyste M&A senior dans ses travaux quotidiens. "
    "LANGUE ET TON : "
    "Tu reponds en francais par defaut, sauf si on te parle dans une autre langue. "
    "Ton registre est professionnel, precis et direct. Tu vas droit au but, sans formules de politesse inutiles. "
    "Tu n inventes jamais une information : si tu n es pas certain, tu le signales explicitement. "
    "SOURCING ET RIGUEUR FACTUELLE : "
    "Des qu un chiffre, une donnee de marche, un multiple boursier ou une information factuelle provient d une source externe, "
    "tu indiques la source entre parentheses immediatement apres : ex. (Source : Bloomberg, avril 2025). "
    "Si tu n as pas acces a une donnee precise, tu le dis clairement et tu proposes une methodologie pour l obtenir. "
    "TYPES DE TACHES : "
    "1. ACTUALITE ET VEILLE : resumes structures avec date, source et impact potentiel sur les transactions. "
    "2. RECHERCHE D ACQUEREURS : classe par categorie (strategiques, PE, family offices) avec rationale et acquisitions recentes. "
    "3. COMPARABLES BOURSIERS : tableau structure Societe | Pays | Capi | VE | EV/EBITDA LTM | EV/EBITDA NTM | EV/CA | P/E avec source et date. "
    "4. RELECTURE : corrige fautes et formulations, retourne le texte corrige avec modifications identifiees si demande. "
    "5. ANALYSE DOCUMENTS : points cles, risques, inconsistances, regard critique niveau banquier senior. "
    "6. ETUDES DE MARCHE : taille, TCAM, acteurs, tendances — chaque chiffre source. "
    "FORMAT : utilise titres, sous-titres et tableaux. "
    "Termine chaque analyse complexe par une section Points cles a retenir de 3 a 5 bullets maximum."
)


# ==========================================
# 4. FONCTIONS
# ==========================================

def extract_pdf(file_bytes):
    try:
        reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception as e:
        return "[Erreur PDF : " + str(e) + "]"


def extract_pptx(file_bytes):
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        out = []
        for i, slide in enumerate(prs.slides, 1):
            parts = ["--- Slide " + str(i) + " ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
            out.append("\n".join(parts))
        return "\n\n".join(out)
    except Exception as e:
        return "[Erreur PPTX : " + str(e) + "]"


def to_b64(file_bytes):
    return base64.b64encode(file_bytes).decode("utf-8")


def build_messages(history, system_prompt):
    msgs = [{"role": "system", "content": system_prompt}]
    for msg in history:
        msgs.append({"role": msg["role"], "content": msg["api_content"]})
    return msgs


def get_system_with_date(base_prompt):
    today = datetime.now().strftime("%A %d %B %Y")
    return "Nous sommes le " + today + ". Tu as acces a internet et aux informations en temps reel. " + base_prompt


def friendly_error(err):
    s = str(err)
    if "429" in s:
        return (
            "**Quota atteint sur ce modele.**\n\n"
            "- Patientez 30 secondes et reessayez\n"
            "- Changez de modele dans la sidebar\n"
            "- Verifiez votre solde sur openrouter.ai"
        )
    if "404" in s:
        return "**Modele introuvable.** L ID a peut-etre change. Essayez un autre modele."
    if "401" in s or "403" in s:
        return "**Erreur d authentification.** Verifiez votre cle API dans secrets.toml"
    return "**Erreur API :** " + s


def generate_image_pollinations(prompt_text, style_suffix, width=1024, height=1024):
    """
    Utilise Pollinations.ai — completement gratuit, sans cle API.
    Retourne les bytes de l image.
    """
    full_prompt = prompt_text
    if style_suffix:
        full_prompt = prompt_text + ", " + style_suffix
    encoded = quote(full_prompt)
    url = (
        "https://image.pollinations.ai/prompt/" + encoded
        + "?width=" + str(width)
        + "&height=" + str(height)
        + "&nologo=true"
        + "&enhance=true"
    )
    resp = requests.get(url, timeout=90)
    resp.raise_for_status()
    return resp.content, url


# ==========================================
# 5. AUTHENTIFICATION
# ==========================================
def check_password():
    if st.session_state.get("authenticated"):
        return True
    st.markdown(
        "<div style='max-width:360px;margin:80px auto;text-align:center;'>"
        "<div style='font-family:\"Cormorant Garamond\",serif;font-size:2rem;color:#ddd8ce;font-weight:600;'>"
        "IA <span style='color:#b8995a;'>CIC</span></div>"
        "<div style='font-size:0.72rem;color:#6b7089;letter-spacing:.14em;text-transform:uppercase;"
        "margin-top:6px;margin-bottom:32px;'>Acces collaborateurs</div></div>",
        unsafe_allow_html=True
    )
    _, col, _ = st.columns([1, 2, 1])
    with col:
        pwd = st.text_input("pwd", type="password", label_visibility="collapsed", placeholder="Mot de passe...")
        if st.button("Acceder", use_container_width=True):
            if pwd == st.secrets.get("APP_PASSWORD", ""):
                st.session_state["authenticated"] = True
                st.rerun()
            else:
                st.error("Mot de passe incorrect")
    return False


if not check_password():
    st.stop()


# ==========================================
# 6. SESSION
# ==========================================
if "messages" not in st.session_state:
    st.session_state.messages = []
if "system_prompt" not in st.session_state:
    st.session_state.system_prompt = DEFAULT_SYSTEM
if "last_loaded_conv" not in st.session_state:
    st.session_state.last_loaded_conv = None
if "active_tab" not in st.session_state:
    st.session_state.active_tab = "chat"


# ==========================================
# 7. CLIENT API
# ==========================================
@st.cache_resource
def get_client():
    return OpenAI(
        base_url="https://openrouter.ai/api/v1",
        api_key=st.secrets["OPENROUTER_API_KEY"],
        default_headers={
            "HTTP-Referer": "https://iacic.streamlit.app",
            "X-Title": "IA CIC",
        }
    )

client = get_client()


# ==========================================
# 8. SIDEBAR
# ==========================================
with st.sidebar:
    st.markdown(
        "<div class='sidebar-header'>"
        "<div class='sidebar-brand'>IA <span>CIC</span></div>"
        "<div class='sidebar-tagline'>Intelligence Artificielle Collaborative</div>"
        "</div>",
        unsafe_allow_html=True
    )

    st.markdown("<div class='section-label'>Modele</div>", unsafe_allow_html=True)
    selected = st.selectbox("Modele", list(MODELS.keys()), label_visibility="collapsed")
    m = MODELS[selected]
    ts = TIER_STYLES[m["tier"]]
    st.markdown(
        "<div style='font-size:0.71rem;color:" + ts["color"] + ";font-weight:600;"
        "letter-spacing:.08em;margin-top:5px;'>● " + ts["label"] + "</div>"
        "<div style='font-size:0.78rem;color:#8b90a8;margin-top:4px;line-height:1.55;'>" + m["desc"] + "</div>"
        "<div style='font-size:0.69rem;color:#4a5270;margin-top:3px;'>💳 " + m["price"] + "</div>",
        unsafe_allow_html=True
    )

    st.markdown("<div class='section-divider'></div>", unsafe_allow_html=True)

    st.markdown("<div class='section-label'>Instructions systeme</div>", unsafe_allow_html=True)
    st.session_state.system_prompt = st.text_area(
        "sys", value=st.session_state.system_prompt,
        height=110, label_visibility="collapsed",
        placeholder="Definir le comportement de l IA..."
    )

    st.markdown("<div class='section-divider'></div>", unsafe_allow_html=True)

    st.markdown("<div class='section-label'>Recherche web</div>", unsafe_allow_html=True)
    web_search = st.toggle("Activer la recherche internet", value=True)
    if web_search:
        st.markdown(
            "<div style='font-size:0.72rem;color:#5a9e7a;margin-top:2px;'>"
            "🌐 Connectee a internet — actu en temps reel</div>",
            unsafe_allow_html=True
        )
    else:
        st.markdown(
            "<div style='font-size:0.72rem;color:#6b7089;margin-top:2px;'>"
            "📴 Hors ligne — connaissance jusqu en 2024</div>",
            unsafe_allow_html=True
        )

    st.markdown("<div class='section-divider'></div>", unsafe_allow_html=True)

    st.markdown("<div class='section-label'>Piece jointe</div>", unsafe_allow_html=True)
    uploaded_file = st.file_uploader(
        "Fichier", type=["pdf", "pptx", "png", "jpg", "jpeg", "webp"],
        label_visibility="collapsed"
    )
    if uploaded_file:
        ext = uploaded_file.name.split(".")[-1].upper()
        name_short = uploaded_file.name[:30] + ("..." if len(uploaded_file.name) > 30 else "")
        st.markdown(
            "<div class='file-chip'>📎 " + name_short +
            "<span style='color:#b8995a;margin-left:4px;'>" + ext + "</span></div>",
            unsafe_allow_html=True
        )

    st.markdown("<div class='section-divider'></div>", unsafe_allow_html=True)

    msg_count = len(st.session_state.messages)
    user_count = sum(1 for x in st.session_state.messages if x["role"] == "user")
    st.markdown(
        "<div class='stat-row'>"
        "<div class='stat-pill'>💬 " + str(msg_count) + " messages</div>"
        "<div class='stat-pill'>👤 " + str(user_count) + " questions</div>"
        "</div>",
        unsafe_allow_html=True
    )

    st.markdown("<div style='height:12px;'></div>", unsafe_allow_html=True)

    # Sauvegarde
    if st.session_state.messages:
        conv_name = datetime.now().strftime("%Y-%m-%d_%H-%M")
        exportable = []
        for msg in st.session_state.messages:
            exportable.append({
                "role": msg["role"],
                "display_content": msg["display_content"],
                "api_content": msg["api_content"] if isinstance(msg["api_content"], str)
                               else msg["display_content"]
            })
        conv_json = json.dumps({
            "name": conv_name,
            "model": selected,
            "system_prompt": st.session_state.system_prompt,
            "messages": exportable
        }, ensure_ascii=False, indent=2)
        st.download_button(
            label="💾  Sauvegarder la conversation",
            data=conv_json.encode("utf-8"),
            file_name="conv_" + conv_name + ".json",
            mime="application/json",
            use_container_width=True
        )

    st.markdown("<div style='height:6px;'></div>", unsafe_allow_html=True)

    st.markdown("<div class='section-label'>Charger une conversation</div>", unsafe_allow_html=True)
    uploaded_conv = st.file_uploader(
        "conv", type=["json"],
        label_visibility="collapsed", key="conv_uploader"
    )
    if uploaded_conv is not None:
        file_id = uploaded_conv.name + "_" + str(uploaded_conv.size)
        if st.session_state.last_loaded_conv != file_id:
            try:
                conv_data = json.loads(uploaded_conv.read().decode("utf-8"))
                restored = []
                for msg in conv_data.get("messages", []):
                    restored.append({
                        "role": msg["role"],
                        "display_content": msg["display_content"],
                        "api_content": msg["api_content"]
                    })
                st.session_state.messages = restored
                if conv_data.get("system_prompt"):
                    st.session_state.system_prompt = conv_data["system_prompt"]
                st.session_state.last_loaded_conv = file_id
                st.success("Conversation restauree : " + conv_data.get("name", ""))
            except Exception as e:
                st.error("Erreur de chargement : " + str(e))

    st.markdown("<div style='height:6px;'></div>", unsafe_allow_html=True)

    if st.button("🗑  Nouvelle conversation"):
        st.session_state.messages = []
        st.session_state.last_loaded_conv = None
        st.rerun()


# ==========================================
# 9. ZONE PRINCIPALE
# ==========================================
tab_chat, tab_image = st.tabs(["💬  Chat", "🎨  Generateur d image"])


# ==========================================
# TAB 1 — CHAT (affichage uniquement)
# ==========================================
with tab_chat:
    model_short = selected.split("—")[0].strip()
    st.markdown(
        "<div class='chat-header'>"
        "<div class='chat-title'>Intelligence <span>Artificielle</span></div>"
        "<div class='chat-meta'>Modele : " + model_short +
        " &nbsp;·&nbsp; " + str(len(st.session_state.messages)) + " messages</div>"
        "</div>",
        unsafe_allow_html=True
    )

    if not st.session_state.messages:
        st.markdown(
            "<div class='welcome-box'>"
            "<h3>Bienvenue sur votre espace IA</h3>"
            "<p>Posez vos questions, analysez vos documents, redigez ou explorez des idees. "
            "Choisissez le modele adapte dans le panneau lateral.</p>"
            "<div class='capability-grid'>"
            "<div class='capability-item'><span class='icon'>📄</span>Analyse PDF et PowerPoint</div>"
            "<div class='capability-item'><span class='icon'>🖼️</span>Lecture d images et captures</div>"
            "<div class='capability-item'><span class='icon'>🔍</span>Etudes de marche et recherches</div>"
            "<div class='capability-item'><span class='icon'>✍️</span>Redaction et synthese</div>"
            "</div></div>",
            unsafe_allow_html=True
        )

    for msg in st.session_state.messages:
        with st.chat_message(msg["role"]):
            st.markdown(msg["display_content"])


# ==========================================
# TAB 2 — GENERATEUR D'IMAGE
# ==========================================
with tab_image:
    st.markdown(
        "<div class='chat-header'>"
        "<div class='chat-title'>Generateur d <span>Images</span></div>"
        "<div class='chat-meta'>Creez des visuels a partir d une description — gratuit, sans limite</div>"
        "</div>",
        unsafe_allow_html=True
    )

    col_left, col_right = st.columns([1, 1], gap="large")

    with col_left:
        st.markdown("<div class='section-label'>Style</div>", unsafe_allow_html=True)
        selected_style = st.selectbox(
            "style", list(IMAGE_STYLES.keys()),
            label_visibility="collapsed"
        )

        st.markdown("<div class='section-label' style='margin-top:14px;'>Format</div>", unsafe_allow_html=True)
        format_choice = st.selectbox(
            "format", ["Carre (1024x1024)", "Paysage (1280x720)", "Portrait (720x1280)"],
            label_visibility="collapsed"
        )
        size_map = {
            "Carre (1024x1024)": (1024, 1024),
            "Paysage (1280x720)": (1280, 720),
            "Portrait (720x1280)": (720, 1280),
        }
        img_w, img_h = size_map[format_choice]

        st.markdown("<div class='section-label' style='margin-top:14px;'>Description</div>", unsafe_allow_html=True)
        img_prompt = st.text_area(
            "img_prompt", height=140,
            label_visibility="collapsed",
            placeholder="Ex : Vue aerienne d une ville europeenne de nuit, architecture haussmannienne, lumières dorées..."
        )
        st.markdown(
            "<div style='font-size:0.72rem;color:#6b7089;margin-bottom:14px;line-height:1.6;'>"
            "💡 Plus votre description est precise (style, couleurs, cadrage, ambiance), meilleur sera le resultat."
            "</div>",
            unsafe_allow_html=True
        )
        generate_btn = st.button("✦  Generer l image", use_container_width=True)

    with col_right:
        st.markdown("<div class='section-label'>Resultat</div>", unsafe_allow_html=True)

        if generate_btn:
            if not img_prompt.strip():
                st.warning("Veuillez entrer une description.")
            else:
                with st.spinner("Generation en cours... (5 a 15 secondes)"):
                    try:
                        style_suffix = IMAGE_STYLES[selected_style]
                        img_bytes, img_url = generate_image_pollinations(
                            img_prompt, style_suffix, img_w, img_h
                        )
                        st.image(img_bytes, use_container_width=True)
                        st.download_button(
                            label="⬇️  Telecharger l image",
                            data=img_bytes,
                            file_name="image_cic_" + datetime.now().strftime("%Y%m%d_%H%M%S") + ".png",
                            mime="image/png",
                            use_container_width=True
                        )
                        st.markdown(
                            "<div style='font-size:0.7rem;color:#4a5270;margin-top:6px;'>Genere via Pollinations.ai</div>",
                            unsafe_allow_html=True
                        )
                    except Exception as e:
                        st.error("Erreur generation image : " + str(e))
        else:
            st.markdown(
                "<div style='background:#1b2030;border:1px dashed #252c3f;border-radius:10px;"
                "height:320px;display:flex;align-items:center;justify-content:center;"
                "color:#3a4060;font-size:2rem;'>✦</div>",
                unsafe_allow_html=True
            )


# ==========================================
# 10. CHAT INPUT — EN DEHORS DES ONGLETS
# pour qu il reste colle en bas de page
# ==========================================
if prompt := st.chat_input("Posez votre question..."):

    api_content = [{"type": "text", "text": prompt}]
    display_content = prompt

    if uploaded_file:
        file_bytes = uploaded_file.read()
        ext = uploaded_file.name.split(".")[-1].lower()

        if ext in ("png", "jpg", "jpeg", "webp"):
            b64 = to_b64(file_bytes)
            mime = "image/jpeg" if ext == "jpg" else "image/" + ext
            api_content.append({
                "type": "image_url",
                "image_url": {"url": "data:" + mime + ";base64," + b64}
            })
            display_content += "\n\n📎 *Image jointe : " + uploaded_file.name + "*"

        elif ext == "pdf":
            text = extract_pdf(file_bytes)
            api_content[0]["text"] += "\n\nContenu du PDF (" + uploaded_file.name + ") :\n" + text
            display_content += "\n\n📎 *PDF joint : " + uploaded_file.name + "*"

        elif ext == "pptx":
            text = extract_pptx(file_bytes)
            api_content[0]["text"] += "\n\nContenu du PowerPoint (" + uploaded_file.name + ") :\n" + text
            display_content += "\n\n📎 *PowerPoint joint : " + uploaded_file.name + "*"

    st.session_state.messages.append({
        "role": "user",
        "api_content": api_content,
        "display_content": display_content,
    })

    # On force le retour sur l onglet chat pour voir la reponse
    with tab_chat:
        with st.chat_message("user"):
            st.markdown(display_content)

        with st.chat_message("assistant"):
            placeholder = st.empty()
            full_response = ""

            try:
                system_with_date = get_system_with_date(st.session_state.system_prompt)
                api_messages = build_messages(st.session_state.messages, system_with_date)
                extra = {}
                if web_search:
                    extra["plugins"] = [{"id": "web"}]
                stream = client.chat.completions.create(
                    model=m["id"],
                    messages=api_messages,
                    stream=True,
                    max_tokens=8192,
                    extra_body=extra if extra else None,
                )
                for chunk in stream:
                    delta = chunk.choices[0].delta
                    if delta and delta.content:
                        full_response += delta.content
                        placeholder.markdown(full_response + "▌")
                placeholder.markdown(full_response)

            except Exception as e:
                full_response = friendly_error(e)
                placeholder.warning(full_response)

    st.session_state.messages.append({
        "role": "assistant",
        "api_content": full_response,
        "display_content": full_response,
    })
